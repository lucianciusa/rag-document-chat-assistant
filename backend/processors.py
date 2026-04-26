import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import os

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv", ".png", ".jpg", ".jpeg", ".bmp"}

def parse_document(file_path: str, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported format '{ext}'. Supported: PDF, DOCX, PPTX, TXT, CSV, MD, PNG, JPG, JPEG, BMP.")
    try:
        if ext == ".pdf":
            return parse_pdf(file_path)
        elif ext in [".md", ".txt", ".csv"]:
            return parse_text(file_path)
        elif ext == ".docx":
            return parse_docx(file_path)
        elif ext == ".pptx":
            return parse_pptx(file_path)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp"]:
            return parse_image(file_path)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Could not parse '{filename}': {e}")

def parse_pdf(file_path: str) -> str:
    try:
        doc = fitz.open(file_path)
    except Exception:
        raise ValueError("PDF file is corrupted or cannot be opened.")
    if doc.is_encrypted:
        raise ValueError("PDF is password-protected and cannot be indexed.")
    text = "".join(page.get_text() + "\n" for page in doc)
    if not text.strip():
        raise ValueError("PDF contains no extractable text (scanned image? Try uploading as PNG/JPG for OCR).")
    return text.strip()

def parse_text(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read().strip()

def parse_docx(file_path: str) -> str:
    try:
        doc = Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs).strip()
        if not text:
            raise ValueError("DOCX contains no extractable text.")
        return text
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"DOCX file is corrupted or invalid: {e}")

def parse_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        if not text.strip():
            raise ValueError("PPTX contains no extractable text.")
        return text.strip()
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"PPTX file is corrupted or invalid: {e}")

def parse_image(file_path: str) -> str:
    try:
        img = Image.open(file_path)
        img.verify()
    except Exception:
        raise ValueError("Image file is corrupted or in an unsupported format.")
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
    except Exception as e:
        raise ValueError(f"OCR failed: {e}")
    if not text.strip():
        raise ValueError("No text detected in image (OCR found nothing).")
    return text.strip()
